"""Knowledge ingestion pipeline — documents → chunks → embeddings → index.

    python ingest.py                 # add anything new, keep what is indexed
    python ingest.py --rebuild       # discard the index and start over
    python ingest.py --dirs a,b      # ingest other directories (also DATA_DIR)

On a large corpus this runs for hours, nearly all of it in the embedding phase,
which reports `embedded X/Y` per batch. Keep that output:

    python -u ingest.py 2>&1 | tee ingest.log

`-u` unbuffers Python, `tee` writes the file and still shows the run. Do NOT
pipe through `tail`/`head` — they hold everything in memory and print only once
the process ends, so a long run looks silent and an interrupted one leaves
nothing at all. The log matters afterwards too: `save_state` writes the manifest
ONCE, at the very end, so an interrupted run leaves no other trace of which
files gave no text and which failed to read.

Written without LangChain, like the agent itself: a PDF/TXT/MD reader, a
recursive splitter, batched OpenAI embeddings, and a FAISS index saved next to
the chunk texts — the same texts BM25 uses at query time.

INCREMENTAL BY CONTENT HASH. Every source file is keyed by the sha256 of its
bytes, so re-running after adding one document embeds that document only. That
is the assignment's "reloads without re-embedding", and the reason a large
private corpus can be grown one folder at a time.
"""

from __future__ import annotations

import argparse
import fnmatch
import hashlib
import json
import os
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path

# `embeddings` first: with the local backend it initialises torch, which must
# precede faiss on macOS (see the note in retriever.py).
import embeddings as emb  # isort: skip
import numpy as np
from pypdf import PdfReader

import ocr
import vectorstore
from config import settings

# Pages recovered by OCR, as (path, page). Recognition makes mistakes, so a
# chunk that came from an image is flagged and the agent is told to treat its
# wording with care.
ocr_pages: set[tuple[str, int]] = set()

TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".rst"}
PDF_SUFFIXES = {".pdf"}
DOCX_SUFFIXES = {".docx"}
XLSX_SUFFIXES = {".xlsx", ".xlsm"}
PPTX_SUFFIXES = {".pptx"}
SUPPORTED = (TEXT_SUFFIXES | PDF_SUFFIXES | DOCX_SUFFIXES
             | XLSX_SUFFIXES | PPTX_SUFFIXES)

INDEX_FILE = "index.faiss"
CHUNKS_FILE = "chunks.json"
MANIFEST_FILE = "manifest.json"


# --------------------------------------------------------------------------- #
# reading
# --------------------------------------------------------------------------- #


def file_digest(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for block in iter(lambda: f.read(1 << 20), b""):
            h.update(block)
    return h.hexdigest()


def read_docx(path: Path) -> str:
    """Paragraphs and table cells of a Word document, in reading order.

    Tables matter more than they look: in the kind of documents people actually
    keep — contracts, invoices, specifications — the numbers live in tables,
    and a reader that takes only paragraphs silently drops exactly the part
    someone will search for. Word has no page concept in the file, so every
    chunk is recorded as page 1.
    """
    from docx import Document

    document = Document(str(path))
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_xlsx(path: Path) -> list[tuple[str, int]]:
    """One entry per sheet; rows serialised as "header: value" pairs.

    A spreadsheet row dumped raw — "1200 | 15.03.2024 | Acme Ltd" — embeds
    badly, because embedding models are trained on prose, not on grids. Written
    out as "Amount: 1200; Date: 15.03.2024; Counterparty: Acme Ltd" the same row
    reads like a sentence and lands in a sensible place in vector space. The
    sheet name goes in too, since "Payments" or "2024 Q3" is often the only
    thing that says what the numbers are.

    Exact figures and document numbers are still mostly found by BM25 — that is
    what lexical search is for. This makes them findable semantically as well.
    """
    from openpyxl import load_workbook

    # read_only: never loads the whole grid; data_only: formulas as values.
    book = load_workbook(str(path), read_only=True, data_only=True)
    out: list[tuple[str, int]] = []
    try:
        for index, sheet in enumerate(book.worksheets, start=1):
            # Excel routinely declares a sheet as 1048575 x 16384 while holding
            # a dozen rows, and read_only mode believes the declaration: one
            # real file here took 142 seconds to cross 200k empty rows, and
            # would have needed ~12 minutes per sheet. reset_dimensions()
            # recomputes the range from the cells that exist — same 14 rows of
            # data, 0.35 seconds.
            if hasattr(sheet, "reset_dimensions"):
                sheet.reset_dimensions()
            lines = [f"Sheet: {sheet.title}"]
            headers: list[str] | None = None
            blank_run = 0
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if c is None else str(c).strip() for c in row]
                if not any(cells):
                    # Second guard, for files where even the recomputed range
                    # is wrong: a long silence means the data has ended.
                    blank_run += 1
                    if blank_run > settings.xlsx_blank_run_limit:
                        break
                    continue
                blank_run = 0
                if headers is None:                 # first non-empty row
                    headers = cells
                    lines.append(" | ".join(h for h in cells if h))
                    continue
                pairs = [f"{h}: {v}" for h, v in zip(headers, cells) if v and h]
                lines.append("; ".join(pairs) if pairs
                             else " | ".join(c for c in cells if c))
            if len(lines) > 1:
                out.append(("\n".join(lines), index))
    finally:
        book.close()
    return out


def read_pptx(path: Path) -> list[tuple[str, int]]:
    """One entry per slide: shapes, tables and the speaker notes.

    The notes are the point. A slide often carries three words and a picture,
    while what the presenter actually meant is written underneath — so a reader
    that takes only the visible text indexes the least informative half.
    """
    from pptx import Presentation

    deck = Presentation(str(path))
    out: list[tuple[str, int]] = []
    for index, slide in enumerate(deck.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Notes: {notes}")
        if parts:
            out.append(("\n".join(parts), index))
    return out


def read_document(path: Path) -> list[tuple[str, int]]:
    """Return [(text, locator)].

    The locator is a page for PDFs, a sheet number for spreadsheets, a slide
    number for decks, and 1 for formats without any structure of their own.
    """
    suffix = path.suffix.lower()

    if suffix in XLSX_SUFFIXES:
        try:
            return read_xlsx(path)
        except Exception as exc:
            print(f"   ! {path.name}: {type(exc).__name__}: {exc}")
            return []

    if suffix in PPTX_SUFFIXES:
        try:
            return read_pptx(path)
        except Exception as exc:
            print(f"   ! {path.name}: {type(exc).__name__}: {exc}")
            return []

    if path.suffix.lower() in DOCX_SUFFIXES:
        try:
            text = read_docx(path)
        except Exception as exc:
            print(f"   ! {path.name}: {type(exc).__name__}: {exc}")
            return []
        return [(text, 1)] if text.strip() else []

    if path.suffix.lower() in PDF_SUFFIXES:
        pages: list[tuple[str, int]] = []
        unreadable: list[int] = []
        try:
            reader = PdfReader(str(path))
        except Exception as exc:          # encrypted, truncated, not really a PDF
            print(f"   ! {path.name}: {type(exc).__name__}: {exc}")
            return []
        for number, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception as exc:              # one broken page is not fatal
                print(f"   ! {path.name} p.{number}: {type(exc).__name__}: {exc}")
                continue
            if len(text.strip()) >= settings.ocr_min_chars:
                pages.append((text, number))
            else:
                unreadable.append(number)

        # Only the pages that gave nothing go to OCR — it costs seconds per
        # page, against milliseconds for a text layer.
        if unreadable and ocr.resolve_backend() != "off":
            print(f"   ↻ {path.name}: {len(unreadable)} page(s) without text → "
                  f"{ocr.describe()}")
            for text, number in ocr.ocr_pdf_pages(path, unreadable):
                pages.append((text, number))
                ocr_pages.add((str(path), number))
        pages.sort(key=lambda item: item[1])
        return pages
    return [(path.read_text(encoding="utf-8", errors="replace"), 1)]


def is_excluded(name: str) -> bool:
    """Does this filename match EXCLUDE_FILES? (fnmatch, case-insensitive)"""
    lowered = name.lower()
    return any(fnmatch.fnmatch(lowered, pattern.strip().lower())
               for pattern in settings.exclude_files.split(",") if pattern.strip())


def discover(dirs: list[str]) -> list[Path]:
    """Find ingestible files, pruning directories that are not a corpus.

    Pointing this at a real folder taught the lesson: a documents directory
    that also holds a git clone contributes hundreds of READMEs and package
    docs from `.venv`, and they surface in search as confident noise. Hidden
    directories and the names in EXCLUDE_DIRS are pruned during the walk, so
    the tree is never even descended.

    EXCLUDE_FILES prunes by filename for a different reason: some files are not
    documents (Office `~$` lock files), and some must never become searchable
    (keys, recovery codes). What was skipped is printed rather than dropped in
    silence — an exclusion that quietly swallows a wanted document is the same
    class of failure as a document that quietly yields no text.
    """
    excluded = {d.strip() for d in settings.exclude_dirs.split(",") if d.strip()}
    found: list[Path] = []
    skipped: list[str] = []
    for raw in dirs:
        root = Path(raw).expanduser()
        if not root.is_absolute():
            root = Path(__file__).parent / root
        if not root.exists():
            print(f"   ! directory not found, skipped: {root}")
            continue
        for current, subdirs, files in os.walk(root):
            subdirs[:] = sorted(d for d in subdirs
                                if d not in excluded and not d.startswith("."))
            for name in sorted(files):
                if Path(name).suffix.lower() not in SUPPORTED or name.startswith("."):
                    continue
                if is_excluded(name):
                    skipped.append(name)
                    continue
                found.append(Path(current) / name)
    if skipped:
        print(f"   ⊘ {len(skipped)} file(s) skipped by EXCLUDE_FILES: "
              + ", ".join(skipped[:6]) + (" …" if len(skipped) > 6 else ""))
    return found


def text_key(text: str) -> str:
    """Identity of a chunk's CONTENT, insensitive to whitespace and case.

    Two exports of one document (.docx and .pdf side by side) rarely differ in
    words but often differ in spacing, so comparing raw strings would call them
    different and index both.
    """
    return hashlib.sha256(" ".join(text.split()).lower().encode("utf-8")).hexdigest()


def collapse_duplicates(pairs: list[tuple[dict, np.ndarray]]
                        ) -> tuple[list[tuple[dict, np.ndarray]], int]:
    """Keep the first chunk of each distinct text, drop the rest.

    Measured on a real corpus: 1637 of 8121 chunks (20%) were exact repeats —
    the same document stored as both .docx and .pdf, "— копия" files, and
    spreadsheets whose own text repeats a row. The cost is not disk. Retrieval
    returned the SAME passage three times for one query, so the three
    RERANK_TOP_N slots that should have carried three facts carried one, and
    the reranker cannot help: to it they are three different chunks.

    First occurrence wins, so an incremental run never re-orders what is
    already indexed. The trade-off, stated plainly: if the file that won is
    later deleted, its text leaves with it even though a copy is still on
    disk — `--rebuild` puts it back.
    """
    seen: set[str] = set()
    kept: list[tuple[dict, np.ndarray]] = []
    for chunk, vector in pairs:
        key = text_key(chunk["text"])
        if key in seen:
            continue
        seen.add(key)
        kept.append((chunk, vector))
    return kept, len(pairs) - len(kept)


# --------------------------------------------------------------------------- #
# chunking
# --------------------------------------------------------------------------- #

SEPARATORS = ["\n\n", "\n", ". ", " ", ""]


def split_text(text: str, size: int, overlap: int) -> list[str]:
    """Recursive character splitter.

    Same idea as RecursiveCharacterTextSplitter: break on the coarsest
    separator that fits — paragraph, line, sentence, word — so a chunk rarely
    ends mid-thought.
    """
    text = text.strip()
    if len(text) <= size:
        return [text] if text else []

    for sep in SEPARATORS:
        if sep == "":                                  # last resort: hard cut
            step = max(1, size - overlap)
            return [text[i:i + size] for i in range(0, len(text), step)]
        parts = text.split(sep)
        if len(parts) == 1:
            continue

        chunks: list[str] = []
        current = ""
        for part in parts:
            candidate = part if not current else current + sep + part
            if len(candidate) <= size:
                current = candidate
                continue
            if current:
                chunks.append(current)
                current = ""
            if len(part) <= size:
                current = part
            else:                       # too long for this level: go finer
                chunks.extend(split_text(part, size, overlap))
        if current:
            chunks.append(current)

        return [c for c in stitch(chunks, overlap) if c.strip()]
    return [text]


def carry_tail(previous: str, overlap: int) -> str:
    """The tail of one chunk, carried into the next — cut at a WORD boundary.

    `previous[-overlap:]` is a raw character slice, so it lands mid-word almost
    every time. Measured on the built index: 70.8% of chunks (4562 of 6448)
    began with a word fragment — "s those stated in clause 4.1…" is the tail of
    "as". The embedder then sees a junk first token, the reranker sees it too,
    and the agent quotes it into the report verbatim.

    A tail with no space in it at all is dropped rather than truncated: one
    unbroken 100-character token carries no context worth keeping.
    """
    if overlap <= 0 or not previous:
        return ""
    tail = previous[-overlap:]
    if len(previous) > overlap:            # the slice actually cut something
        space = tail.find(" ")
        tail = tail[space + 1:] if space != -1 else ""
    return tail.strip()


def stitch(chunks: list[str], overlap: int) -> list[str]:
    """Carry each chunk's tail into the next, so a fact sitting on a boundary
    stays retrievable from either side."""
    if overlap <= 0 or len(chunks) < 2:
        return chunks
    out = [chunks[0]]
    for previous, chunk in zip(chunks, chunks[1:]):
        tail = carry_tail(previous, overlap)
        out.append(f"{tail} {chunk}".strip() if tail else chunk)
    return out


def document_metadata(path: Path) -> dict:
    """Дата документа и — обязательно — её происхождение.

    Дат две, и они отвечают на разные вопросы. `created` внутри самого файла —
    когда документ создал автор; mtime — когда файл лёг на этот диск. Замер по
    корпусу: внутренняя дата есть у 239 из 287 PDF и DOCX (83%), у остальных
    остаётся только mtime.

    Поэтому пишется не только дата, но и `date_source`. Без него через месяц
    никто не отличит настоящую дату подписания договора от даты, когда папку
    скопировали с другого ноутбука, — а выглядеть они будут одинаково
    убедительно.
    """
    created = None
    try:
        if path.suffix.lower() in PDF_SUFFIXES:
            info = PdfReader(str(path)).metadata
            created = info.creation_date if info else None
        elif path.suffix.lower() in DOCX_SUFFIXES | XLSX_SUFFIXES | PPTX_SUFFIXES:
            # docx, xlsx и pptx — это OOXML, то есть zip с одинаковым
            # docProps/core.xml. Читать его напрямую дешевле, чем поднимать три
            # разные библиотеки ради одного поля.
            import re
            import zipfile
            with zipfile.ZipFile(path) as archive:
                core = archive.read("docProps/core.xml").decode("utf-8", "replace")
            found = re.search(r"<dcterms:created[^>]*>([^<]+)<", core)
            created = datetime.fromisoformat(
                found.group(1).replace("Z", "+00:00")) if found else None
    except Exception:
        # Битый или нестандартный контейнер — не повод терять файл: ниже есть
        # mtime, который есть всегда.
        created = None

    if created:
        return {"date": created.date().isoformat(), "date_source": "document"}
    stamp = datetime.fromtimestamp(path.stat().st_mtime)
    return {"date": stamp.date().isoformat(), "date_source": "filesystem"}


_mail_metadata: dict[str, dict] | None = None


def mail_metadata(path: Path) -> dict:
    """Метаданные письма для файла, если он пришёл вложением.

    Вложение — обычный документ по содержимому и почтовый объект по
    происхождению. Без второй половины его нельзя найти вопросом «что мне
    присылал X»: фильтр участников смотрит на поля письма, а у PDF их нет.
    Справочник строится один раз за прогон и только если такие файлы вообще
    встретились — корпус документов не обязан знать о существовании почты.
    """
    global _mail_metadata

    attachments_root = (Path(__file__).parent / settings.mail_attachments_dir).resolve()
    try:
        relative = path.resolve().relative_to(attachments_root)
    except ValueError:
        return {}

    if _mail_metadata is None:
        try:
            from mailprep.attachments import metadata_by_folder
            _mail_metadata = metadata_by_folder(
                Path(__file__).parent / settings.mail_db)
        except Exception as exc:            # почта не должна ронять индексацию
            print(f"   ! mail metadata unavailable: {type(exc).__name__}: {exc}")
            _mail_metadata = {}

    # Первый сегмент пути внутри mail/attachments — каталог письма.
    return _mail_metadata.get(relative.parts[0], {}) if relative.parts else {}


def chunk_metadata(path: Path) -> dict:
    """Всё, что известно о файле помимо его текста.

    Приоритет дат: собственная дата документа, затем дата письма, и только
    потом mtime. Порядок выведен из промаха, а не из вкуса: сначала дата письма
    стояла первой, и устав 2020 года выводился как «2026-06-16», потому что
    письмо с ним переслали в этом году. Зато когда своей даты у файла нет,
    дата письма всё равно лучше mtime — mtime у вложения это момент, когда его
    скачал `imap_fetch`, то есть чистый артефакт нашего же конвейера.
    """
    metadata = {**document_metadata(path), **mail_metadata(path)}
    if metadata.get("mail_date") and metadata["date_source"] == "filesystem":
        metadata["date"] = metadata["mail_date"]
        metadata["date_source"] = "mail"
    return metadata


def chunk_document(path: Path, digest: str) -> list[dict]:
    chunks = []
    metadata = chunk_metadata(path)
    for text, page in read_document(path):
        pieces = split_text(text, settings.chunk_size, settings.chunk_overlap)
        for position, body in enumerate(pieces):
            chunk = {
                "id": hashlib.sha256(
                    f"{digest}:{page}:{position}".encode()).hexdigest()[:16],
                "text": body,
                "source": path.name,
                "path": str(path),
                "page": page,
                **metadata,
            }
            if (str(path), page) in ocr_pages:
                chunk["ocr"] = True
            chunks.append(chunk)
    return chunks


# --------------------------------------------------------------------------- #
# embedding
# --------------------------------------------------------------------------- #


def embed(texts: list[str]) -> np.ndarray:
    """Passages, not queries — the distinction matters for prefix-trained models."""
    return emb.embed_texts(texts, is_query=False, progress=True)


# --------------------------------------------------------------------------- #
# index i/o
# --------------------------------------------------------------------------- #


def index_dir() -> Path:
    path = Path(__file__).parent / settings.index_dir
    path.mkdir(parents=True, exist_ok=True)
    return path


def load_chunks() -> list[dict]:
    path = index_dir() / CHUNKS_FILE
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else []


def load_manifest() -> tuple[dict, dict | None]:
    """Return (files: path -> digest, embedding signature or None).

    The manifest records WHICH MODEL built the index, not just which files went
    into it. Without that, changing EMBEDDING_MODEL in .env and forgetting to
    rebuild produces an index that still answers — with neighbours computed in
    a different vector space. No exception, no warning, just quietly wrong
    retrieval. With it, the next query fails loudly and says what to do.
    """
    path = index_dir() / MANIFEST_FILE
    if not path.exists():
        return {}, None
    data = json.loads(path.read_text(encoding="utf-8"))
    if "files" not in data:                       # pre-signature manifest
        return data, None
    return data.get("files", {}), data.get("embedding")


def load_vectors() -> np.ndarray | None:
    store = vectorstore.get_store()
    return store.all_vectors() if store.exists() else None


def save_state(chunks: list[dict], files: dict, vectors: np.ndarray,
               barren: list[str] = None) -> None:
    directory = index_dir()
    manifest = {
        "embedding": {**emb.signature(), "dim": int(vectors.shape[1])},
        "files": files,
        # Kept so the list survives the run: these are the documents a later
        # OCR pass has to target.
        "no_text": sorted(barren or []),
    }
    vectorstore.get_store().write(vectors, chunks)
    (directory / CHUNKS_FILE).write_text(
        json.dumps(chunks, ensure_ascii=False), encoding="utf-8")
    (directory / MANIFEST_FILE).write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


# --------------------------------------------------------------------------- #
# pipeline
# --------------------------------------------------------------------------- #


def ingest(dirs: list[str] | None = None, rebuild: bool = False) -> dict:
    # Embedding a corpus with a local model is the memory-hungriest thing this
    # project does; check before spending minutes on it, not after.
    import preflight

    preflight.guard(includes_reranker=False)   # ingestion never loads it

    dirs = dirs or settings.ingest_dirs
    print(f"scanning: {', '.join(dirs)}")
    files = discover(dirs)
    if not files:
        print("nothing to ingest.")
        return {"files": 0, "chunks": 0, "embedded": 0}

    if rebuild:
        chunks, manifest, vectors = [], {}, None
    else:
        chunks, vectors = load_chunks(), load_vectors()
        manifest, stored_embedding = load_manifest()
        if vectors is not None and len(vectors) != len(chunks):
            print("   ! index and chunk list disagree — rebuilding from scratch")
            chunks, manifest, vectors = [], {}, None
        elif vectors is not None:
            # Appending vectors from a second model into one index would make
            # the two halves incomparable, and nothing downstream could tell.
            try:
                emb.check_compatible(stored_embedding)
            except RuntimeError as exc:
                print(f"\n{exc}\n")
                raise SystemExit(2)

    # chunk id -> vector, so surviving chunks keep their embedding
    known: dict[str, np.ndarray] = {}
    if vectors is not None:
        known = {c["id"]: v for c, v in zip(chunks, vectors)}

    kept = list(chunks)
    fresh: list[dict] = []
    barren: list[str] = []          # files that produced no text at all
    for path in files:
        digest = file_digest(path)
        if manifest.get(str(path)) == digest:
            print(f" = {path.name} (unchanged)")
            continue
        produced = chunk_document(path, digest)
        if not produced:
            # A file that yields nothing is the quietest failure in the whole
            # pipeline: it is recorded as processed, contributes no chunks, and
            # the agent later answers "not in the knowledge base" about a
            # document that is sitting right there. Measured on a real corpus:
            # ~15% of PDFs were scans with no text layer at all.
            print(f" ∅ {path.name} — no text extracted")
            barren.append(str(path))
        else:
            print(f" + {path.name}")
        kept = [c for c in kept if c["path"] != str(path)]   # replace, never duplicate
        fresh.extend(produced)
        manifest[str(path)] = digest

    for known_path in list(manifest):                        # deleted on disk
        if not Path(known_path).exists():
            print(f" - {Path(known_path).name} (removed)")
            kept = [c for c in kept if c["path"] != known_path]
            manifest.pop(known_path)

    if not fresh and vectors is not None and len(kept) == len(chunks):
        print(f"index up to date: {len(chunks)} chunks from {len(manifest)} files")
        return {"files": len(manifest), "chunks": len(chunks), "embedded": 0}

    fresh_vectors = None
    if fresh:
        # Drop repeats BEFORE paying for them. collapse_duplicates() below is
        # the safety net and runs over everything, but by then the vectors are
        # bought: on this corpus 20% of chunks are duplicates, which is ~50
        # minutes of embedding computed and immediately thrown away.
        seen = {text_key(chunk["text"]) for chunk in kept if chunk["id"] in known}
        unique: list[dict] = []
        for chunk in fresh:
            key = text_key(chunk["text"])
            if key in seen:
                continue
            seen.add(key)
            unique.append(chunk)
        if len(unique) < len(fresh):
            print(f"   ⊙ {len(fresh) - len(unique)} duplicate chunk(s) skipped "
                  "before embedding")
            fresh = unique
        print(f"embedding {len(fresh)} new chunk(s) with "
              f"{emb.describe(emb.signature())}")
        fresh_vectors = embed([c["text"] for c in fresh])

    # Chunks and vectors are carried as PAIRS from here on. They used to be two
    # lists filtered by different conditions, which is how a chunk without a
    # vector could stay in chunks.json and shift every later row by one — the
    # index then answers with somebody else's text, confidently and silently.
    paired: list[tuple[dict, np.ndarray]] = [
        (chunk, known[chunk["id"]]) for chunk in kept if chunk["id"] in known
    ]
    if fresh_vectors is not None:
        paired.extend(zip(fresh, fresh_vectors))

    paired, dropped = collapse_duplicates(paired)
    if dropped:
        print(f"   ⊙ {dropped} duplicate chunk(s) collapsed "
              "(same text already in the index)")

    if not paired:
        print("nothing indexed.")
        return {"files": len(manifest), "chunks": 0, "embedded": 0}

    final_chunks = [chunk for chunk, _ in paired]
    final_vectors = np.stack([vector for _, vector in paired]).astype("float32")
    save_state(final_chunks, manifest, final_vectors, barren)
    print(f"\nindex saved to {index_dir()}")
    print(f"  files: {len(manifest)} | chunks: {len(final_chunks)} | "
          f"newly embedded: {len(fresh)} | dim: {final_vectors.shape[1]}")

    if barren:
        print(f"\n⚠️  {len(barren)} file(s) produced NO text and are in the "
              "index in name only:")
        for path_str in barren[:10]:
            print(f"     ∅ {Path(path_str).name}")
        if len(barren) > 10:
            print(f"     … and {len(barren) - 10} more (full list in "
                  f"{MANIFEST_FILE} under \"no_text\")")
        print("   Most often these are scanned PDFs — images with no text "
              "layer. They need OCR; until then the agent will honestly say "
              "it knows nothing about them.")
    return {"files": len(manifest), "chunks": len(final_chunks),
            "embedded": len(fresh)}


def relabel() -> dict:
    """Re-derive chunk metadata in place, for every chunk already indexed.

    Metadata is not part of the vector. A contract's text does not change
    because we now also know its date or who sent it, so re-embedding 7347
    chunks to add a field would buy nothing and cost hours on this machine. The
    chunks keep their ids, their order and their vectors; only the payload
    grows.

    This exists because ingestion is incremental by file digest: the files are
    unchanged, so an ordinary run correctly skips all of them and any newly
    added field would never appear on what is already indexed.

    Files that have since been deleted keep whatever metadata they had — the
    chunk stays in the index either way, and inventing a date for a file nobody
    can read is worse than leaving the old one.
    """
    directory = index_dir()
    chunks, vectors = load_chunks(), load_vectors()
    if not chunks or vectors is None:
        print(f"no index in {directory} — nothing to relabel.")
        return {"chunks": 0, "labelled": 0}
    if len(chunks) != len(vectors):
        print("   ! index and chunk list disagree — run --rebuild instead.")
        raise SystemExit(2)

    raw = json.loads((directory / MANIFEST_FILE).read_text(encoding="utf-8"))
    # Метаданные читаются по файлу, а чанков у файла десятки: кэш превращает
    # 7347 чтений с диска в 430.
    by_path: dict[str, dict] = {}
    labelled = mailed = missing = 0
    sources = Counter()
    for chunk in chunks:
        path = Path(chunk["path"])
        if chunk["path"] not in by_path:
            if path.exists():
                by_path[chunk["path"]] = chunk_metadata(path)
            else:
                by_path[chunk["path"]] = {}
                missing += 1
        metadata = by_path[chunk["path"]]
        if metadata:
            chunk.update(metadata)
            labelled += 1
            sources[metadata["date_source"]] += 1
            mailed += bool(metadata.get("sender_email"))

    save_state(chunks, raw.get("files", {}),
               np.asarray(vectors, dtype="float32"), raw.get("no_text", []))
    print(f"relabelled {directory}")
    print(f"  chunks: {labelled} of {len(chunks)} | files read: {len(by_path)}")
    print("  date source: " + ", ".join(f"{name} {count}"
                                        for name, count in sources.most_common()))
    print(f"  chunks with mail metadata: {mailed}")
    if missing:
        print(f"  ! {missing} file(s) no longer on disk — metadata left as it was")
    return {"chunks": len(chunks), "labelled": labelled}


def clean() -> dict:
    """Apply the current EXCLUDE_FILES and de-duplication policy in place.

    Nothing is read from the source documents and nothing is embedded: every
    surviving chunk already has its vector in the index, and dropping rows from
    a flat index is just keeping the rows that stay. That matters when the
    corpus took four hours to embed and the reason to clean it — a secret that
    should never have been indexed — is urgent.

    Not a substitute for the policy at ingest time: this removes what is
    already there, `discover()` prevents it from arriving.
    """
    directory = index_dir()
    chunks, vectors = load_chunks(), load_vectors()
    if not chunks or vectors is None:
        print(f"no index in {directory} — nothing to clean.")
        return {"files": 0, "chunks": 0, "removed": 0}
    if len(chunks) != len(vectors):
        print("   ! index and chunk list disagree — run --rebuild instead.")
        raise SystemExit(2)

    raw = json.loads((directory / MANIFEST_FILE).read_text(encoding="utf-8"))
    manifest, barren = raw.get("files", {}), raw.get("no_text", [])
    before = len(chunks)

    def name_of(chunk: dict) -> str:
        return chunk.get("source") or Path(chunk["path"]).name

    excluded_sources = sorted({name_of(c) for c in chunks if is_excluded(name_of(c))})
    paired = [(c, v) for c, v in zip(chunks, vectors) if not is_excluded(name_of(c))]
    by_policy = before - len(paired)

    paired, duplicates = collapse_duplicates(paired)

    # A file barred by policy must also leave the manifest, or the next run
    # sees it as "already indexed" and never notices it is gone.
    manifest = {p: d for p, d in manifest.items() if not is_excluded(Path(p).name)}
    barren = [p for p in barren if not is_excluded(Path(p).name)]

    if not paired:
        print("everything was removed — refusing to write an empty index.")
        raise SystemExit(2)

    final_chunks = [c for c, _ in paired]
    final_vectors = np.stack([v for _, v in paired]).astype("float32")
    save_state(final_chunks, manifest, final_vectors, barren)

    print(f"cleaned {directory}")
    if excluded_sources:
        print(f"  ⊘ removed by EXCLUDE_FILES ({by_policy} chunk(s)): "
              + ", ".join(excluded_sources))
    print(f"  ⊙ duplicate chunks collapsed: {duplicates}")
    print(f"  chunks: {before} → {len(final_chunks)} "
          f"(−{before - len(final_chunks)}, "
          f"{100 * (before - len(final_chunks)) / before:.1f}%)")
    print(f"  files in manifest: {len(manifest)}")
    return {"files": len(manifest), "chunks": len(final_chunks),
            "removed": before - len(final_chunks)}


def main() -> None:
    parser = argparse.ArgumentParser(description="Build the knowledge index.")
    parser.add_argument("--dirs", help="comma-separated directories to ingest")
    parser.add_argument("--rebuild", action="store_true",
                        help="discard the existing index and embed everything")
    parser.add_argument("--migrate-store", metavar="FROM:TO",
                        help="move vectors between engines without re-embedding, "
                             "e.g. --migrate-store faiss:qdrant")
    parser.add_argument("--clean", action="store_true",
                        help="apply EXCLUDE_FILES and de-duplication to the "
                             "existing index without re-embedding anything")
    parser.add_argument("--relabel", action="store_true",
                        help="attach mail metadata (sender, recipients, date) "
                             "to attachment chunks already indexed, without "
                             "re-embedding anything")
    args = parser.parse_args()
    dirs = [d.strip() for d in args.dirs.split(",")] if args.dirs else None
    try:
        if args.migrate_store:
            source, _, target = args.migrate_store.partition(":")
            moved = vectorstore.migrate(source, target)
            print(f"moved {moved} vectors {source} → {target} in {index_dir()}")
            print(f"set VECTOR_BACKEND={target} in .env to use it")
        elif args.clean:
            clean()
        elif args.relabel:
            relabel()
        else:
            ingest(dirs, rebuild=args.rebuild)
    except KeyboardInterrupt:
        print("\ninterrupted; the previous index is untouched.")
        sys.exit(130)


if __name__ == "__main__":
    main()
